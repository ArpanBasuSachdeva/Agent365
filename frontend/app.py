import streamlit as st
import requests
import io
import os
import tempfile
from pathlib import Path
import zipfile
from datetime import datetime

# For file previews
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Page configuration
st.set_page_config(
    layout="wide",
    page_title="Agent365",
    page_icon="📄",
    initial_sidebar_state="expanded"
)

# Custom CSS for modern, clean light theme UI
st.markdown("""
    <style>
    /* Force light theme */
    .stApp {
        background-color: #ffffff;
    }
    .main .block-container {
        background-color: #ffffff;
    }
    [data-testid="stSidebar"] {
        background-color: #f9fafb;
    }
    
    /* Main header */
    .main-header {
        font-size: 2.5rem;
        font-weight: 700;
        color: #1f2937;
        margin-bottom: 0.5rem;
    }
    .subtitle {
        font-size: 1.1rem;
        color: #6b7280;
        margin-bottom: 2rem;
    }
    .footer {
        text-align: center;
        color: #9ca3af;
        font-size: 0.9rem;
        margin-top: 3rem;
        padding-top: 2rem;
        border-top: 1px solid #e5e7eb;
    }
    
    /* Buttons */
    .stButton>button {
        width: 100%;
        background-color: #3b82f6;
        color: white;
        font-weight: 600;
        padding: 0.75rem 1.5rem;
        border-radius: 0.5rem;
        border: none;
    }
    .stButton>button:hover {
        background-color: #2563eb;
    }
    
    /* Info boxes */
    .success-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #d1fae5;
        border-left: 4px solid #10b981;
        margin: 1rem 0;
    }
    .error-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #fee2e2;
        border-left: 4px solid #ef4444;
        margin: 1rem 0;
    }
    .info-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #dbeafe;
        border-left: 4px solid #3b82f6;
        margin: 1rem 0;
    }
    
    /* Text colors for light theme */
    h1, h2, h3, h4, h5, h6 {
        color: #1f2937;
    }
    p, div, span {
        color: #1f2937;
    }
    </style>
""", unsafe_allow_html=True)

# Initialize session state
if 'processed_file' not in st.session_state:
    st.session_state.processed_file = None
if 'file_name' not in st.session_state:
    st.session_state.file_name = None
if 'task_summary' not in st.session_state:
    st.session_state.task_summary = None
if 'processing_history' not in st.session_state:
    st.session_state.processing_history = []

# Backend URL
BACKEND_URL = os.getenv("BACKEND_URL", "http://localhost:8000/process")
BACKEND_USERNAME = os.getenv("BACKEND_USERNAME", "admin")
BACKEND_PASSWORD = os.getenv("BACKEND_PASSWORD", "change-me")

# Sidebar
with st.sidebar:
    st.title("📋 Agent365")
    st.markdown("---")
    
    st.subheader("⚙️ Settings")
    backend_url = st.text_input("Backend URL", value=BACKEND_URL, help="FastAPI backend endpoint")
    username = st.text_input("Username", value=BACKEND_USERNAME, help="Basic Auth username")
    password = st.text_input("Password", value=BACKEND_PASSWORD, type="password", help="Basic Auth password")
    
    st.markdown("---")
    
    st.subheader("📊 History")
    if st.session_state.processing_history:
        for i, entry in enumerate(reversed(st.session_state.processing_history[-10:])):
            with st.expander(f"📄 {entry.get('file_name', 'Unknown')} - {entry.get('timestamp', '')}"):
                st.write(f"**Prompt:** {entry.get('prompt', 'N/A')}")
                st.write(f"**Status:** {entry.get('status', 'N/A')}")
                if entry.get('summary'):
                    st.write(f"**Summary:** {entry.get('summary', 'N/A')}")
    else:
        st.info("No processing history yet.")

# Main content
st.markdown('<div class="main-header">Agent365</div>', unsafe_allow_html=True)
st.markdown('<div class="subtitle">AI Office File Editor • Powered by LLMs (Executor + Validator Agents)</div>', unsafe_allow_html=True)

# Main container
col1, col2 = st.columns([1, 1])

with col1:
    st.subheader("📤 Upload File")
    uploaded_file = st.file_uploader(
        "Choose an Office file",
        type=['docx', 'xlsx', 'pptx'],
        help="Upload a .docx, .xlsx, or .pptx file"
    )
    
    if uploaded_file:
        st.info(f"✅ **File uploaded:** {uploaded_file.name} ({uploaded_file.size:,} bytes)")

with col2:
    st.subheader("✍️ Enter Prompt")
    prompt = st.text_area(
        "Describe what you want to do with the file",
        height=150,
        placeholder="e.g., 'Add a summary table at the end', 'Update all dates to 2024', 'Change the title to Annual Report'",
        help="Enter a natural language description of the changes you want"
    )

st.markdown("---")

# Process button
col_btn1, col_btn2, col_btn3 = st.columns([1, 2, 1])
with col_btn2:
    process_button = st.button("🚀 Process File", type="primary", use_container_width=True)

# Processing logic
if process_button:
    if not uploaded_file:
        st.error("❌ Please upload a file first.")
    elif not prompt or not prompt.strip():
        st.error("❌ Please enter a prompt describing what you want to do.")
    else:
        with st.spinner("🔄 Processing your file... This may take a few moments."):
            try:
                # Prepare the file for upload
                file_bytes = uploaded_file.read()
                files = {"file": (uploaded_file.name, file_bytes, uploaded_file.type)}
                data = {"prompt": prompt}
                
                # Make request to backend
                response = requests.post(
                    backend_url,
                    files=files,
                    data=data,
                    auth=(username, password),
                    timeout=300  # 5 minute timeout for processing
                )
                
                # Handle response
                if response.status_code == 200:
                    # Check if response is a file or JSON
                    content_type = response.headers.get('content-type', '')
                    
                    if 'application/json' in content_type:
                        # JSON response (error or info)
                        result = response.json()
                        st.error(f"❌ Error: {result.get('error', 'Unknown error')}")
                        st.session_state.processing_history.append({
                            'file_name': uploaded_file.name,
                            'prompt': prompt,
                            'status': 'Error',
                            'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                        })
                    else:
                        # File response
                        st.success("✅ File processed successfully!")
                        
                        # Get task summary from headers if available
                        task_summary = response.headers.get('X-Task-Summary', 'Task completed successfully')
                        st.session_state.task_summary = task_summary
                        
                        # Store processed file
                        st.session_state.processed_file = response.content
                        st.session_state.file_name = f"updated_{uploaded_file.name}"
                        
                        # Add to history
                        st.session_state.processing_history.append({
                            'file_name': uploaded_file.name,
                            'prompt': prompt,
                            'status': 'Success',
                            'summary': task_summary,
                            'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                        })
                        
                        st.info(f"📝 **Task Summary:** {task_summary}")
                        
                        # Show preview
                        st.markdown("---")
                        st.subheader("👁️ Preview")
                        show_preview(uploaded_file.name, response.content)
                        
                elif response.status_code == 401:
                    st.error("❌ Authentication failed. Please check your username and password in the sidebar.")
                elif response.status_code == 404:
                    st.error(f"❌ Endpoint not found. Please check the backend URL: {backend_url}")
                else:
                    error_msg = f"❌ Error {response.status_code}: {response.text}"
                    st.error(error_msg)
                    st.session_state.processing_history.append({
                        'file_name': uploaded_file.name,
                        'prompt': prompt,
                        'status': f'Error {response.status_code}',
                        'timestamp': pd.Timestamp.now().strftime("%Y-%m-%d %H:%M:%S") if PANDAS_AVAILABLE else "N/A"
                    })
                    
            except requests.exceptions.Timeout:
                st.error("❌ Request timed out. The file processing is taking longer than expected.")
            except requests.exceptions.ConnectionError:
                st.error(f"❌ Could not connect to backend at {backend_url}. Please ensure the FastAPI server is running.")
            except Exception as e:
                st.error(f"❌ An unexpected error occurred: {str(e)}")

# Download section
if st.session_state.processed_file:
    st.markdown("---")
    st.subheader("📥 Download Edited File")
    
    col_dl1, col_dl2, col_dl3 = st.columns([1, 2, 1])
    with col_dl2:
        st.download_button(
            label="⬇️ Download Edited File",
            data=st.session_state.processed_file,
            file_name=st.session_state.file_name,
            mime="application/octet-stream",
            use_container_width=True
        )

# Preview function
def show_preview(file_name, file_content):
    """Show preview of the processed file based on file type"""
    file_ext = Path(file_name).suffix.lower()
    
    try:
        if file_ext == '.docx' and DOCX_AVAILABLE:
            # Preview Word document
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                tmp_file.write(file_content)
                tmp_path = tmp_file.name
            
            try:
                doc = Document(tmp_path)
                st.write("**Document Preview:**")
                
                # Show paragraphs
                for i, para in enumerate(doc.paragraphs[:20]):  # Limit to first 20 paragraphs
                    if para.text.strip():
                        st.markdown(f"**Paragraph {i+1}:** {para.text}")
                
                if len(doc.paragraphs) > 20:
                    st.info(f"... and {len(doc.paragraphs) - 20} more paragraphs")
                
                # Show tables
                if doc.tables:
                    st.write("**Tables:**")
                    for i, table in enumerate(doc.tables):
                        st.write(f"**Table {i+1}:**")
                        table_data = []
                        for row in table.rows:
                            table_data.append([cell.text for cell in row.cells])
                        if table_data:
                            st.table(table_data[:10])  # Show first 10 rows
            finally:
                os.unlink(tmp_path)
                
        elif file_ext == '.xlsx' and PANDAS_AVAILABLE:
            # Preview Excel file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
                tmp_file.write(file_content)
                tmp_path = tmp_file.name
            
            try:
                excel_file = pd.ExcelFile(tmp_path)
                st.write("**Workbook Preview:**")
                
                # Show all sheets
                for sheet_name in excel_file.sheet_names:
                    st.write(f"**Sheet: {sheet_name}**")
                    df = pd.read_excel(tmp_path, sheet_name=sheet_name, nrows=20)  # First 20 rows
                    st.dataframe(df, use_container_width=True)
                    if len(pd.read_excel(tmp_path, sheet_name=sheet_name)) > 20:
                        st.info(f"... and more rows in this sheet")
                    st.markdown("---")
            finally:
                os.unlink(tmp_path)
                
        elif file_ext == '.pptx' and PPTX_AVAILABLE:
            # Preview PowerPoint
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                tmp_file.write(file_content)
                tmp_path = tmp_file.name
            
            try:
                prs = Presentation(tmp_path)
                st.write("**Presentation Preview:**")
                
                for i, slide in enumerate(prs.slides, 1):
                    st.write(f"**Slide {i}:**")
                    # Extract text from slide
                    text_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_content.append(shape.text.strip())
                    
                    if text_content:
                        for text in text_content[:5]:  # First 5 text elements
                            st.write(f"  • {text}")
                    else:
                        st.write("  (No text content)")
                    st.markdown("---")
            finally:
                os.unlink(tmp_path)
        else:
            st.info("📄 Preview not available for this file type. Please download the file to view it.")
            if not DOCX_AVAILABLE and file_ext == '.docx':
                st.warning("💡 Install `python-docx` for Word document preview: `pip install python-docx`")
            if not PANDAS_AVAILABLE and file_ext == '.xlsx':
                st.warning("💡 Install `pandas` and `openpyxl` for Excel preview: `pip install pandas openpyxl`")
            if not PPTX_AVAILABLE and file_ext == '.pptx':
                st.warning("💡 Install `python-pptx` for PowerPoint preview: `pip install python-pptx`")
                
    except Exception as e:
        st.warning(f"⚠️ Could not generate preview: {str(e)}")
        st.info("📄 Please download the file to view it.")

# Footer
st.markdown("---")
st.markdown(
    '<div class="footer">Secure • Agentic Loop • Reliable Outputs</div>',
    unsafe_allow_html=True
)

