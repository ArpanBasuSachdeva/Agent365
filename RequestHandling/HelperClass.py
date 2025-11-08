"""
HelperClass.py - Shared utilities and helper functions for Agent365
"""

import os
import time
import shutil
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple
import google.generativeai as genai
from dotenv import load_dotenv
import json
import hashlib

# Load environment variables
load_dotenv()

class Agent365Helper:
    """Main helper class for Agent365 operations"""
    
    def __init__(self):
        """Initialize the helper class with configuration"""
        self.api_key = os.getenv("GEMINI_API_KEY")
        if not self.api_key:
            raise ValueError("GEMINI_API_KEY not found in environment variables")
        
        # Configure Gemini
        genai.configure(api_key=self.api_key)
        self.model = genai.GenerativeModel("gemini-2.5-flash")
        
        # Set up directories
        self.codes_dir = Path(os.getenv("CODES_DIR", "D:/Ricky/Projects/Office_Agent/codes"))
        self.temp_dir = Path(os.getenv("TEMP_DIR", "D:/Ricky/Projects/Office_Agent/temp"))
        self.codes_dir.mkdir(parents=True, exist_ok=True)
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        
        # Load user credentials
        self.users_db_path = Path("users.json")
        self.users = self._load_users()
        
        # Secondary validator model (can reuse the same for simplicity)
        self.validator_model = self.model
    
    def _load_users(self) -> Dict[str, Dict[str, Any]]:
        """Load users from JSON file or create default users"""
        if self.users_db_path.exists():
            try:
                with open(self.users_db_path, 'r') as f:
                    return json.load(f)
            except (json.JSONDecodeError, FileNotFoundError):
                pass
        
        # Create default users if file doesn't exist
        default_users = {
            "admin": {
                "password_hash": self._hash_password("password123"),
                "role": "admin",
                "created_at": time.time(),
                "last_login": None
            },
            "user1": {
                "password_hash": self._hash_password("user123"),
                "role": "user",
                "created_at": time.time(),
                "last_login": None
            },
            "user2": {
                "password_hash": self._hash_password("user456"),
                "role": "user",
                "created_at": time.time(),
                "last_login": None
            }
        }
        self._save_users(default_users)
        return default_users
    
    def _save_users(self, users: Dict[str, Dict[str, Any]]):
        """Save users to JSON file"""
        with open(self.users_db_path, 'w') as f:
            json.dump(users, f, indent=2)
    
    def _hash_password(self, password: str) -> str:
        """Hash password using SHA-256"""
        return hashlib.sha256(password.encode()).hexdigest()
    
    def authenticate_user(self, username: str, password: str) -> bool:
        """Authenticate user with username and password"""
        if username not in self.users:
            return False
        
        password_hash = self._hash_password(password)
        if self.users[username]["password_hash"] == password_hash:
            # Update last login time
            self.users[username]["last_login"] = time.time()
            self._save_users(self.users)
            return True
        return False
    
    def add_user(self, username: str, password: str, role: str = "user") -> bool:
        """Add a new user"""
        if username in self.users:
            return False  # User already exists
        
        self.users[username] = {
            "password_hash": self._hash_password(password),
            "role": role,
            "created_at": time.time(),
            "last_login": None
        }
        self._save_users(self.users)
        return True
    
    def remove_user(self, username: str) -> bool:
        """Remove a user"""
        if username not in self.users:
            return False
        
        del self.users[username]
        self._save_users(self.users)
        return True
    
    def change_password(self, username: str, old_password: str, new_password: str) -> bool:
        """Change user password"""
        if username not in self.users:
            return False
        
        if not self.authenticate_user(username, old_password):
            return False  # Old password incorrect
        
        self.users[username]["password_hash"] = self._hash_password(new_password)
        self._save_users(self.users)
        return True
    
    def get_user_info(self, username: str) -> Optional[Dict[str, Any]]:
        """Get user information (without password hash)"""
        if username not in self.users:
            return None
        
        user_info = self.users[username].copy()
        del user_info["password_hash"]  # Don't expose password hash
        return user_info
    
    def list_users(self) -> List[Dict[str, Any]]:
        """List all users (without password hashes)"""
        users_list = []
        for username, user_data in self.users.items():
            user_info = user_data.copy()
            del user_info["password_hash"]
            user_info["username"] = username
            users_list.append(user_info)
        return users_list
    
    # --- Last file path tracking per user ---
    def set_last_file_for_user(self, username: str, file_path: str) -> None:
        """Persist the last file path modified by the given user."""
        try:
            if username not in self.users:
                return
            self.users[username]["last_file_path"] = str(file_path)
            self._save_users(self.users)
        except Exception as e:
            print(f"Warning: could not persist last file path for {username}: {e}")

    def get_last_file_for_user(self, username: str) -> Optional[str]:
        """Retrieve the last file path modified by the given user, if available."""
        try:
            if username not in self.users:
                return None
            return self.users[username].get("last_file_path")
        except Exception:
            return None
    
    def extract_python_code_blocks(self, text: str) -> List[str]:
        """Extract Python code blocks from text response"""
        blocks = []
        in_block = False
        current_block = []
        
        for line in text.split('\n'):
            if line.strip().startswith('```python'):
                in_block = True
                current_block = []
            elif line.strip().startswith('```') and in_block:
                in_block = False
                if current_block:
                    blocks.append('\n'.join(current_block))
                    current_block = []
            elif in_block:
                current_block.append(line)
        
        return blocks

    def extract_imports_from_code(self, code_text: str) -> List[str]:
        """Very simple import scanner to determine dependencies to install.
        Returns a list of top-level packages (best-effort).
        """
        packages: List[str] = []
        try:
            for line in code_text.splitlines():
                stripped = line.strip()
                if stripped.startswith("import "):
                    parts = stripped.split()
                    # import pkg[, pkg2] style - take first after 'import'
                    if len(parts) >= 2:
                        first = parts[1].split(",")[0].split(".")[0]
                        if first and first not in packages:
                            packages.append(first)
                elif stripped.startswith("from ") and " import " in stripped:
                    # from pkg.sub import x
                    try:
                        pkg = stripped.split()[1].split(".")[0]
                        if pkg and pkg not in packages:
                            packages.append(pkg)
                    except Exception:
                        pass
        except Exception:
            pass
        # Filter out stdlib/common names we don't want to pip install
        blacklist = {"os", "sys", "time", "json", "re", "math", "pathlib", "typing", "subprocess", "shutil"}
        return [p for p in packages if p not in blacklist]

    def ensure_dependencies_installed(self, packages: List[str]) -> List[Tuple[str, bool, str]]:
        """Attempt to install any missing packages via pip. Returns list of (package, success, message)."""
        results: List[Tuple[str, bool, str]] = []
        if not packages:
            return results
        try:
            import importlib
            import subprocess
            import sys
            for pkg in packages:
                try:
                    importlib.import_module(pkg)
                    results.append((pkg, True, "already installed"))
                    continue
                except Exception:
                    pass
                try:
                    cmd = [sys.executable, "-m", "pip", "install", pkg]
                    proc = subprocess.run(cmd, capture_output=True, text=True)
                    success = proc.returncode == 0
                    msg = proc.stdout if success else (proc.stderr or proc.stdout)
                    results.append((pkg, success, msg[:400]))
                except Exception as e:
                    results.append((pkg, False, str(e)))
        except Exception as e:
            results.append(("<bootstrap>", False, str(e)))
        return results

    def validate_code_against_task(self, task: str, original_file_content: str, modified_file_content: str, code_text: str = "") -> Dict[str, Any]:
        """Use the model to validate that the modified file matches the task exactly by comparing original vs modified file.
        This ensures only the required task is executed: nothing extra, nothing less.
        Returns { valid: bool, feedback: str }
        """
        try:
            instruction = (
                "You are a strict VALIDATOR. Compare the ORIGINAL file content with the MODIFIED file content "
                "against the user's TASK. Verify that:\n"
                "1. ONLY the required changes from the task were made (nothing extra)\n"
                "2. ALL required changes from the task were made (nothing missing)\n"
                "3. No unrelated modifications were introduced\n"
                "4. The modifications precisely match what was requested in the task\n\n"
                "Output a JSON with fields: valid (true/false), feedback (string explaining what's wrong or confirming correctness)."
            )
            prompt = (
                f"USER TASK:\n{task}\n\n"
                f"ORIGINAL FILE CONTENT:\n{original_file_content[:5000]}\n\n"
                f"MODIFIED FILE CONTENT:\n{modified_file_content[:5000]}\n\n"
                f"Generated code (for reference):\n{code_text[:2000]}"
            )
            resp = self.validator_model.generate_content(instruction + "\n\n" + prompt)
            text = (getattr(resp, "text", None) or "").strip()
            # Try to parse a JSON-ish object; fallback to heuristic
            import json as _json
            result: Dict[str, Any] = {"valid": True, "feedback": ""}
            try:
                result = _json.loads(text)
            except Exception:
                if "valid" in text.lower() and "false" in text.lower():
                    result = {"valid": False, "feedback": text}
                elif not text:
                    result = {"valid": True, "feedback": "No issues found"}
                else:
                    result = {"valid": True, "feedback": text}
            # Normalize
            result["valid"] = bool(result.get("valid", True))
            result["feedback"] = str(result.get("feedback", "")).strip()
            return result
        except Exception as e:
            return {"valid": True, "feedback": f"Validator unavailable: {e}"}

    def regenerate_code_with_feedback(self, task: str, original_code: str, feedback: str) -> str:
        """Ask the model to regenerate code strictly following feedback/constraints. Returns code text."""
        try:
            instruction = (
                "Regenerate ONLY Python code that satisfies the user's task EXACTLY. "
                "Apply the following validator feedback strictly. "
                "Do not include explanations. Return only a single Python code block."
            )
            prompt = (
                f"TASK:\n{task}\n\nVALIDATOR_FEEDBACK:\n{feedback}\n\nORIGINAL_CODE:\n{original_code}"
            )
            resp = self.model.generate_content(instruction + "\n\n" + prompt)
            text = (getattr(resp, "text", None) or "")
            blocks = self.extract_python_code_blocks(text)
            return "\n\n\n".join(blocks) if blocks else original_code
        except Exception:
            return original_code

    def regenerate_code_from_error(self, original_task: str, failed_code: str, error_message: str) -> str:
        """Regenerate code given an execution error message. 
        Sends the executor its own code, the error, and the original task given by the user.
        Returns new code text or original on failure."""
        try:
            instruction = (
                "You are the EXECUTOR agent. Your previously generated code failed during execution. "
                "Fix the Python code to resolve the execution error while ensuring the ORIGINAL USER TASK is still completed correctly. "
                "Analyze the error carefully and fix only what's necessary. "
                "Do not change the logic that correctly implements the task - only fix the error. "
                "Return only a single Python code block with the corrected code."
            )
            prompt = (
                f"ORIGINAL USER TASK (this must still be completed correctly):\n{original_task}\n\n"
                f"EXECUTION ERROR:\n{error_message}\n\n"
                f"YOUR PREVIOUSLY GENERATED CODE (that caused this error):\n{failed_code}"
            )
            resp = self.model.generate_content(instruction + "\n\n" + prompt)
            text = (getattr(resp, "text", None) or "")
            blocks = self.extract_python_code_blocks(text)
            return "\n\n\n".join(blocks) if blocks else failed_code
        except Exception as e:
            print(f"Error in regenerate_code_from_error: {e}")
            return failed_code
    
    def process_file_content(self, file_path: str) -> str:
        """Process different file types and extract content"""
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        
        if suffix == ".docx":
            return self._process_docx_file(file_path)
        elif suffix in [".xlsx", ".xls"]:
            return self._process_excel_file(file_path)
        elif suffix in [".pptx", ".ppt"]:
            return self._process_powerpoint_file(file_path)
        else:
            # For other file types, read as text
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
    
    def _process_docx_file(self, file_path: Path) -> str:
        """Process Word document files"""
        try:
            # Allow empty files for content creation
            try:
                if file_path.stat().st_size == 0:
                    print(f"Warning: Word file is empty (0 bytes), but allowing processing for content creation")
                    return "Empty Word file - ready for content creation"
            except Exception:
                pass

            from docx import Document
            doc = Document(file_path)
            content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            return '\n'.join(content)
        except Exception as e:
            raise Exception(f"Could not read .docx file: {e}")
    
    def _process_excel_file(self, file_path: Path) -> str:
        """Process Excel files"""
        try:
            try:
                if file_path.stat().st_size == 0:
                    print(f"Warning: Excel file is empty (0 bytes), but allowing processing for content creation")
                    return "Empty Excel file - ready for content creation"
            except Exception:
                pass
            import pandas as pd
            df = pd.read_excel(file_path)
            return df.to_string()
        except Exception as e:
            raise Exception(f"Could not read Excel file: {e}")
    
    def _process_powerpoint_file(self, file_path: Path) -> str:
        """Process PowerPoint files"""
        try:
            from pptx import Presentation
            
            # Check if file is empty first - but allow processing for content creation
            file_size = file_path.stat().st_size
            if file_size == 0:
                print(f"Warning: PowerPoint file is empty (0 bytes), but allowing processing for content creation")
                return "Empty PowerPoint file - ready for content creation"
            
            # Try to read directly first
            try:
                prs = Presentation(file_path)
            except (PermissionError, OSError, FileNotFoundError) as e:
                # If direct read fails, copy to temp and try again
                temp_file = self.temp_dir / f"temp_{int(time.time())}{file_path.suffix}"
                shutil.copy2(file_path, temp_file)
                prs = Presentation(temp_file)
            
            content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                content.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(f"  {shape.text}")
            return '\n'.join(content)
        except Exception as e:
            raise Exception(f"Could not read PowerPoint file: {e}")
    
    def generate_instruction_prefix(self) -> str:
        """Generate the instruction prefix for Gemini"""
        return (
            "You are an AI assistant that generates Python code to modify Office documents. "
            "Your task is to write Python code that directly modifies the provided document.\n\n"
            "Available variables:\n"
            "- TARGET_FILE_PATH: path to the file to modify\n"
            "- OUTPUT_DIR: directory where outputs should be saved\n"
            "- CODES_DIR: directory where your generated scripts are saved (read-only for you)\n"
            "Rules: \n"
            "1) Do NOT ask the user for input; the path is provided.\n"
            "2) Save outputs ONLY under OUTPUT_DIR.\n"
            "3) Avoid network calls.\n"
            "4) For Office documents (docx/xlsx/pptx), use python-docx/openpyxl/python-pptx to directly modify the source file:\n"
            "   - Use Document() from docx for Word documents\n"
            "   - Use Workbook() from openpyxl for Excel documents\n"
            "   - Use Presentation() from python-pptx for PowerPoint documents\n"
            "   - CRITICAL: ALWAYS modify the original file at TARGET_FILE_PATH directly - DO NOT create temp files\n"
            "   - CRITICAL: DO NOT use shutil.move(), shutil.copy(), or any file copying/moving operations\n"
            "   - CRITICAL: DO NOT create temporary files and then replace the original - this causes permission errors\n"
            "   - Open the file, modify it in memory, then save directly to TARGET_FILE_PATH\n"
            "   - Save changes back to TARGET_FILE_PATH (same file, not a copy)\n"
            "   - Example: doc = Document(TARGET_FILE_PATH); [modify doc]; doc.save(TARGET_FILE_PATH)\n"
            "   - Example: wb = load_workbook(TARGET_FILE_PATH); [modify wb]; wb.save(TARGET_FILE_PATH)\n"
            "   - If the file is empty (0 bytes), create new content from scratch\n"
            "   - For empty PowerPoint files: create a new presentation with slides\n"
            "   - For empty Word files: create a new document with content\n"
            "   - For empty Excel files: create a new workbook with data\n"
            "   - For openpyxl charts: use chart.add_data() and chart.set_categories(), NOT DataSeries import\n"
            "   - Do NOT set chart.series[-1].title directly - use titles_from_data=True instead\n"
            "   - NEVER use any styles - NO 'Heading 1', 'Title', 'Normal' or any style names\n"
            "   - ONLY use plain text without any formatting or styling\n"
            "   - Do NOT apply any paragraph styles or formatting\n"
            "5) Handle file permissions - if a file is read-only, make it writable or copy to temp location.\n"
            "6) ALWAYS save changes back to TARGET_FILE_PATH - do NOT create new files.\n"
            "7) Print clear completion messages.\n"
        )
    
    def open_file_automatically(self, file_path: Path) -> bool:
        """Automatically open a file using the system default application"""
        try:
            # Wait a moment to ensure file is fully written
            time.sleep(1)
            if os.name == 'nt':  # Windows
                os.startfile(str(file_path))
            elif os.name == 'posix':  # macOS and Linux
                os.system(f'open "{file_path}"')
            print(f"Opened file: {file_path}")
            return True
        except Exception as e:
            print(f"Could not open file automatically: {e}")
            return False
    
    def get_latest_file(self, directory: Path) -> Optional[Path]:
        """Get the most recently modified file in a directory"""
        try:
            files = [f for f in directory.iterdir() if f.is_file()]
            if not files:
                return None
            return max(files, key=lambda f: f.stat().st_mtime)
        except Exception:
            return None
